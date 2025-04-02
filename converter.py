import os
import sys
import subprocess
import argparse
import requests
import json
import base64
from PIL import Image
from io import BytesIO
from dotenv import load_dotenv

load_dotenv()

def process_files(input_data, output_data):
    """
    入力データからファイルを処理し、結果を出力データに格納します。
    
    Args:
        input_data (dict): 入力データを含む辞書
        output_data (dict): 出力データを格納する辞書
    """
    # input_dataからファイルパスを取得
    files = input_data.get('files', [])
    if not files:
        output_data['errors']['general'] = "入力ファイルが指定されていません。"
        return

    for file_path in files:
        if not os.path.exists(file_path):
            output_data['errors'][file_path] = "ファイルが見つかりません"
            continue
            
        markdown_content = convert_to_markdown(file_path)
        if markdown_content:
            output_data['results'][file_path] = markdown_content
        else:
            output_data['errors'][file_path] = "変換に失敗しました"

def convert_to_markdown(input_file, system_prompt=None):
    """
    入力ファイルの形式に応じて適切なプロンプトを選択し、
    Gemini APIを使用してMarkdown形式に変換します。
    """

    file_extension = os.path.splitext(input_file)[1].lower()

    if file_extension == '.pdf':
        system_prompt = """
        あなたは、PDFファイルの内容を説明する専門家です。以下の指示に従って、PDFファイルの内容を説明するMarkdown形式の文書を作成してください。

        1.  画像から読み取れるテキスト情報を抽出してください。
        2.  画像内の図や表について、その内容を文章で説明してください。図や表がない場合は説明しないでください。
        3.  推測や推論は行わず、画像から直接読み取れる情報のみを記述してください。
        4.  画像のタイトルをそのまま使用し、その後に説明を続けてください。タイトルがない場合は説明しないでください。
        5.  Markdown形式: 説明はMarkdown形式で記述してください。
        """
        text_content = ""
    elif file_extension in ['.jpg', '.jpeg', '.png', '.JPG', '.PNG']:
        system_prompt = """
        あなたは、画像ファイルの内容を説明する専門家です。以下の指示に従って、画像の内容を説明するMarkdown形式の文書を作成してください。

        1.  画像の説明: 画像の内容を詳細に説明してください。
        2.  Markdown形式: 説明はMarkdown形式で記述してください。
        3.  画像のファイル名: 画像のファイル名を記述してください。
        4.  その他:
            *   必要に応じて、注釈やコメントを追加してください。
        """
        text_content = ""
    elif file_extension == '.ppt' or file_extension == '.pptx':
        system_prompt = """
        あなたは、PPTファイルをMarkdown形式に変換する専門家です。以下の指示に従って、PPTファイルの内容を説明するMarkdown形式で出力してください。

        1.  スライドの構造: 各スライドをMarkdownの見出し（`#`、`##`など）で表現してください。
        2.  タイトル: スライドのタイトルを最も大きな見出し（`#`）として記述してください。
        3.  リード文: スライドのリード文（概要や導入文）がある場合は、タイトルの直後に記述してください。
        4.  ボディ: スライドの本文の内容を記述してください。箇条書き、段落、表など、適切なMarkdown形式を使用してください。
        5.  図表の記述: スライド内の図や表については、可能な限り内容を理解し、Markdown形式で説明を記述してください。
            *   図の場合: 図の簡単な説明と、図が示している内容を要約してください。必要に応じて、代替テキスト（alt text）を記述してください。
            *   表の場合: 表の構造をMarkdownのテーブル形式で再現し、各行と列の内容を記述してください。難しい場合は、表の内容を箇条書きで記述してください。
        6.  その他:
            *   スライドのノート（発表者ノート）がある場合は、本文の後に注釈として記述してください。
            *   重要なキーワードやフレーズは**強調**してください。

        出力例:

        ```markdown
        # スライドのタイトル

        このスライドのリード文です。

        - ポイント1
        - ポイント2
        - ポイント3

        ![図の説明](image.png)

        表の説明:

        | 列1 | 列2 |
        | --- | --- |
        | 値1 | 値2 |
        | 値3 | 値4 |

        発表者ノート: このスライドでは、〇〇について説明します。
        ```
        """
        text_content = ""
    else:
        print("サポートされていないファイル形式です。PDF、PPT、または画像ファイルを入力してください。")
        return None

    # ファイルの内容をテキストとして抽出
    try:
        if file_extension == '.pdf':
            text_content, image_data = extract_text_and_images_from_pdf(input_file)
        elif file_extension == '.ppt' or file_extension == '.pptx':
            text_content, image_data = extract_text_and_images_from_pptx(input_file)
        elif file_extension in ['.jpg', '.jpeg', '.png', '.JPG', '.PNG']:
            text_content, image_data = extract_text_and_images_from_image(input_file)  # 画像ファイルなので画像データは空
        else:
            return None
    except Exception as e:
        print(f"ファイルの内容抽出中にエラーが発生しました: {e}")
        return None

    # image_dataがNoneの場合の処理を追加
    if image_data is None:
        print("画像データが抽出できませんでした。テキストのみで処理を続行します。")
        image_data = []  # 空のリストで初期化

    # Gemini APIを呼び出してMarkdown形式に変換
    markdown_content = call_gemini_api(system_prompt, text_content, image_data)

    return markdown_content

def extract_text_and_images_from_pdf(pdf_file):
    """PDFファイルからテキストと画像を抽出します (pymupdfを使用)。"""
    try:
        import fitz  # pymupdf

        pdf_document = fitz.open(pdf_file)
        text_content = ""
        image_data = []

        for page_number in range(pdf_document.page_count):
            page = pdf_document.load_page(page_number)
            text_content += page.get_text() + "\n"  # 各ページのテキストを抽出

            # 画像の抽出
            images = page.get_images(full=True)
            for img_index, img in enumerate(images):
                xref = img[0]
                base_image = pdf_document.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                image_data.append({
                    "mime_type": f"image/{image_ext}",
                    "data": image_bytes
                })

        pdf_document.close()
        return text_content, image_data

    except Exception as e:
        print(f"PDFからのテキストと画像の抽出中にエラーが発生しました: {e}")
        return None, None


def extract_text_and_images_from_image(input_file):
    """imageファイルから画像を抽出する処理"""
    try:
        img = Image.open(input_file)
        print(f"Image format: {img.format}")
        if img.format == 'MPO':
            img = img.convert('RGB')
            format = 'JPEG'
        else:
            format = img.format
        print(f"Image format after conversion: {format}")
        buffered = BytesIO()
        img.save(buffered, format=format)
        img_bytes = buffered.getvalue()
        img_str = base64.b64encode(img_bytes).decode('utf-8')
        mime_type = f"image/{format.lower()}"
        image_data = [{'mime_type': mime_type, 'data': img_str}]
        text_content = ""
        return text_content, image_data
    except Exception as e:
        print(f"画像ファイルの処理中にエラーが発生しました: {e}, file: {input_file}")
        return None, None

    #OCRする必要がないと見てコメントアウト
    # def extract_text_from_image(image_file):
    #     """画像ファイルからテキストを抽出します (pytesseractを使用)。"""
    #     try:
    #         import pytesseract
    #         from PIL import Image

    #         # Tesseract OCRのパスを設定 (環境に合わせて変更)
    #         # pytesseract.pytesseract.tesseract_cmd = r'/usr/local/bin/tesseract'  # macOSの場合

    #         img = Image.open(image_file)
        #         text = pytesseract.image_to_string(img, lang='jpn')  # 日本語でOCR
        #         return text
    #     except Exception as e:
        #         print(f"画像からのテキスト抽出中にエラーが発生しました: {e}")
        #         return None

def extract_text_and_images_from_pptx(pptx_file):
    """PPTXファイルからテキストと画像を抽出します。"""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(pptx_file)
        text_content = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
        image_data = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_data.append({"mime_type": shape.image.content_type, "data": shape.image.blob})

        return text_content, image_data
    except Exception as e:
        print(f"PPTXからのテキストと画像の抽出中にエラーが発生しました: {e}")
        return None, None

def call_gemini_api(system_prompt, text_content, image_data):
    """
    Gemini APIを呼び出してMarkdown形式に変換します。
    """
    gemini_api_key = os.environ.get("gemini_api_key")
    if not gemini_api_key:
        print("Gemini APIキーが設定されていません。")
        return None

    url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateContent?key={gemini_api_key}"
    headers = {
        'Content-Type': 'application/json'
    }

    contents = []
    # システムプロンプトを追加
    contents.append({"role": "user", "parts": [{"text": system_prompt}]})

    # テキストコンテンツを追加
    contents.append({"role": "user", "parts": [{"text": text_content}]})

   # 画像データを追加
    for image in image_data:
        if isinstance(image["data"], str):
            # Already base64 encoded
            img_str = image["data"]
        else:
            img_str = base64.b64encode(image["data"]).decode('utf-8')
        contents.append({
            "role": "user",
            "parts": [{
                "inline_data": {
                    "mime_type": image["mime_type"],
                    "data": img_str
                }
            }]
        })

    data = {
        "contents": contents
    }

    try:
        response = requests.post(url, headers=headers, data=json.dumps(data))
        response.raise_for_status()  # エラーレスポンスをチェック
        response_json = response.json()
        markdown_content = response_json['candidates'][0]['content']['parts'][0]['text']
        return markdown_content
    except requests.exceptions.RequestException as e:
        print(f"Gemini APIへのリクエスト中にエラーが発生しました: {e}")
        return None
    except KeyError as e:
        print(f"Gemini APIレスポンスに予期しないキーがありませんでした: {e}")
        return None

if __name__ == "__main__":
    # 標準入力からJSONデータを読み取る
    try:
        input_json = sys.stdin.read()
        input_data = json.loads(input_json)
    except json.JSONDecodeError:
        input_data = {
            'files': sys.argv[1:] if len(sys.argv) > 1 else []
        }
    except Exception as e:
        print(f"入力データの読み取り中にエラーが発生しました: {e}")
        sys.exit(1)

    output_data = {
        'results': {},
        'errors': {}
    }
    
    # ファイルの処理
    process_files(input_data, output_data)
    
    # output_dataディレクトリの作成
    output_dir = 'output_data'
    os.makedirs(output_dir, exist_ok=True)
    
    # 結果をJSONファイルとして保存
    output_file = os.path.join(output_dir, 'output_data.json')
    try:
        with open(output_file, 'w', encoding='utf-8') as f:
            json.dump(output_data, f, ensure_ascii=False, indent=2)
        print(f"\n結果を {output_file} に保存しました。")
    except Exception as e:
        print(f"結果の保存中にエラーが発生しました: {e}")
    
    # 各ファイルの結果を個別のテキストファイルとして保存
    for file_path, content in output_data['results'].items():
        # 元のファイル名から拡張子を除き、.txtを付加
        base_name = os.path.splitext(os.path.basename(file_path))[0]
        output_text_file = os.path.join(output_dir, f"{base_name}.txt")
        try:
            with open(output_text_file, 'w', encoding='utf-8') as f:
                f.write(content)
            print(f"変換結果を {output_text_file} に保存しました。")
        except Exception as e:
            print(f"{output_text_file} の保存中にエラーが発生しました: {e}")
    
    # 結果の表示（デバッグ用）
    if output_data['errors'].get('general'):
        print(f"\n=== エラー ===")
        print(output_data['errors']['general'])
    
    for file_path, content in output_data['results'].items():
        print(f"\n=== {file_path} ===")
        print(content)
    
    for file_path, error in output_data['errors'].items():
        if file_path != 'general':  # 一般エラーは既に表示済み
            print(f"\n=== {file_path} のエラー ===")
            print(error)
